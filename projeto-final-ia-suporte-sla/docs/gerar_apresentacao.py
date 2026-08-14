# -*- coding: utf-8 -*-
"""Gera a apresentacao (pptx) conforme item 7.2 do enunciado."""
import sys, pathlib, json
sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[1]))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config.settings import EVID_DIR, ROOT

M = json.load(open(EVID_DIR / "metrics.json", encoding="utf-8"))
K = json.load(open(ROOT/"data"/"curated"/"dashboard"/"dashboard_data.json", encoding="utf-8"))["kpis"]

NAVY = RGBColor(0x1E,0x27,0x61); TEAL = RGBColor(0x02,0x83,0x90)
INK = RGBColor(0x22,0x2B,0x3A); GRAY = RGBColor(0x5b,0x67,0x70)
LIGHT = RGBColor(0xF2,0xF4,0xF8); WHITE = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.background.fill; r.solid(); r.fore_color.rgb = bg
    return s

def box(s, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill if fill else WHITE
    if line: sp.line.color.rgb = line; sp.line.width = Pt(1.2)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, lines, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Calibri", space=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str): lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def title(s, t, sub=None):
    txt(s, 0.6, 0.45, 12.1, 0.9, t, size=30, color=NAVY, bold=True, font="Cambria")
    if sub: txt(s, 0.62, 1.25, 12.1, 0.5, sub, size=14, color=TEAL, bold=True)

def pic(s, path, x, y, w=None, h=None):
    if pathlib.Path(path).exists():
        kw = {}
        if w: kw["width"] = Inches(w)
        if h: kw["height"] = Inches(h)
        return s.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)

# ---- 1. CAPA (dark) ----
s = slide(NAVY)
txt(s, 0.9, 2.0, 11.5, 1.8, "Pipeline de Dados em Nuvem para Aprendizagem de Máquina",
    size=40, color=WHITE, bold=True, font="Cambria")
txt(s, 0.95, 3.7, 11.3, 0.7, "Previsão de violação de SLA em chamados de suporte técnico",
    size=19, color=RGBColor(0xCA,0xDC,0xFC))
txt(s, 0.95, 4.85, 11.5, 1.7, [
    "Grupo: João Paulo A. S. Monteiro, Rogério B. dos Anjos e",
    "Pedro Felipe de Moraes Carrijo",
    "IFG - Pós-Graduação em IA Aplicada - Módulo 2",
    "Profs. Sirlon Diniz - Raphael Gomes - Otávio Calaça   |   Goiânia, jun/2026",
], size=13, color=RGBColor(0xCA,0xDC,0xFC), space=6)

# ---- 2. VISAO GERAL ----
s = slide(); title(s, "Visão geral da solução", "Do dado bruto à decisão do gestor de suporte")
txt(s, 0.6, 1.9, 6.3, 4.6, [
    "Problema: antecipar quais chamados vão violar o SLA de resolução.",
    "",
    "• Domínio: suporte técnico (service desk SaaS)",
    "• Decisor: gestor da equipe de suporte",
    "• Decisão: priorizar a fila e realocar agentes antes de o prazo estourar",
    "• Tarefa de ML: classificação binária (sla_breach)",
], size=15.5, space=8)
for i,(v,l) in enumerate([(f"{int(K['total_chamados']):,}".replace(',','.'),"chamados"),
                          (f"{K['taxa_violacao_pct']:.0f}%","violam SLA"),
                          (f"{K['csat_medio']:.1f}","CSAT médio")]):
    box(s, 7.4+i*1.85, 2.2, 1.7, 1.6, fill=LIGHT)
    txt(s, 7.4+i*1.85, 2.45, 1.7, 0.8, v, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 7.4+i*1.85, 3.25, 1.7, 0.4, l, size=12, color=GRAY, align=PP_ALIGN.CENTER)
box(s, 7.4, 4.1, 5.3, 2.2, fill=NAVY)
txt(s, 7.7, 4.34, 4.8, 0.5, "Stack do projeto", size=15, bold=True, color=WHITE)
txt(s, 7.7, 4.95, 4.8, 1.2, [
    "• Python, Airflow, dbt",
    "• Snowflake / DuckDB",
    "• AWS (S3, Lambda, CloudFormation), Metabase",
], size=13.5, color=RGBColor(0xCA,0xDC,0xFC), space=8)

# ---- 3. DADOS ----
s = slide(); title(s, "Dados utilizados e atributos extraídos", "3 fontes: estruturada + não estruturada + semiestruturada")
rows = [["Fonte","Tipo","Conteúdo"],
        ["tickets.csv","Estruturado","categoria, prioridade, canal, tier, tempos, CSAT"],
        ["ticket_messages.jsonl","Não estruturado","texto livre do cliente"],
        ["ticket_events.json","Semiestruturado","trilha de eventos/status (JSON)"]]
tb = s.shapes.add_table(4,3, Inches(0.6),Inches(1.95), Inches(7.0),Inches(2.2)).table
for c in range(3): tb.columns[c].width = Inches([2.3,1.7,3.0][c])
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.text=val
        pr=cell.text_frame.paragraphs[0]; pr.runs[0].font.size=Pt(11)
        pr.runs[0].font.bold = ri==0
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if ri==0 else (LIGHT if ri%2 else WHITE)
        if ri==0: pr.runs[0].font.color.rgb=WHITE
txt(s, 0.6, 4.4, 7.0, 2.4, [
    "Atributos extraídos:",
    "• Texto: comprimento, nº palavras, urgência (léxico), flags (reembolso, cancelamento, erro)",
    "• Logs: nº de eventos, mudanças de status, reabertura",
    "• Estruturado: limpeza, padronização, tratamento de ausentes, partes de data",
], size=13.5, space=6)
box(s, 8.0, 1.95, 4.7, 4.6, fill=LIGHT)
txt(s, 8.3, 2.25, 4.2, 0.5, "Sem vazamento de dados", size=16, bold=True, color=TEAL)
txt(s, 8.3, 2.95, 4.2, 3.4, [
    "A decisão ocorre logo após a 1ª resposta.",
    "Por isso NÃO usamos como atributo:",
    "• o tempo de resolução (desfecho)",
    "• o CSAT (só existe no fim)",
    "",
    "Resultado: tabela mart_ml_features com 25 atributos + alvo.",
], size=13.5, space=8, color=INK)

# ---- 4. ARQUITETURA ----
s = slide(); title(s, "Arquitetura da solução", "Equivalente 100% em serviços AWS (item 4.5)")
pic(s, ROOT/"infra"/"arquitetura_aws.png", 1.87, 1.85, h=5.4)

# ---- 5. PIPELINE ELT ----
s = slide(); title(s, "Pipeline de ELT", "Airflow orquestra; dbt transforma no Snowflake")
etapas = ["Ingestão","Processamento\n+ atributos","Treino e\navaliação ML","Carga no\nwarehouse","dbt build\n+ testes","Export\ndashboard"]
for i,e in enumerate(etapas):
    box(s, 0.6+i*2.05, 2.1, 1.8, 1.1, fill=LIGHT, line=TEAL)
    txt(s, 0.6+i*2.05, 2.2, 1.8, 0.95, e.split("\n"), size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=0)
    if i<5: txt(s, 0.6+i*2.05+1.78, 2.35, 0.3, 0.6, ">", size=20, color=TEAL, bold=True)
txt(s, 0.6, 3.7, 6.0, 3.2, [
    "Camadas dbt:",
    "• Staging: stg_tickets, stg_text_features, stg_log_features, stg_predictions",
    "• Dimensões: dim_channel, dim_product, dim_category, dim_date",
    "• Fatos/marts: fct_tickets, mart_ml_features, fct_predictions, mart_ticket_decision",
], size=13.5, space=6)
box(s, 7.0, 3.8, 5.7, 2.6, fill=NAVY)
txt(s, 7.3, 4.02, 5.2, 0.5, "Qualidade de dados (dbt)", size=15, bold=True, color=RGBColor(0x02,0xC3,0x9A))
txt(s, 7.3, 4.64, 5.2, 1.7, [
    "• 10 testes: not_null, unique, accepted_values, relationships",
    "• Última execução: PASS=22, 0 erros",
    "• Reproduzível: run_pipeline.py roda tudo em ~5s",
], size=13.5, color=WHITE, space=8)

# ---- 6. RESULTADOS ML ----
s = slide(); title(s, "Resultados do modelo de ML", "Baselines, KNN (hard-code + biblioteca) e MLP")
mm = M["metricas"]
def r(k,n):
    d=mm[k]
    return [n] + [f"{d[c]:.3f}".replace(".", ",")
                  for c in ('accuracy','precision','recall','f1','roc_auc')]
data=[["Modelo","Acc","Prec","Recall","F1","AUC"],
      r("baseline_regra_prioridade","Baseline (regra)"),
      r("knn_hardcode","KNN (na mão)"),
      r("knn_sklearn","KNN (sklearn)"),
      r("mlp","MLP *")]
tb=s.shapes.add_table(5,6, Inches(0.6),Inches(1.95), Inches(6.6),Inches(2.4)).table
for ci,wd in enumerate([2.1,0.9,0.9,0.9,0.9,0.9]): tb.columns[ci].width=Inches(wd)
for ri,row in enumerate(data):
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.text=val; pr=cell.text_frame.paragraphs[0]
        pr.runs[0].font.size=Pt(11); pr.runs[0].font.bold = ri==0 or ri==4
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if ri==0 else (RGBColor(0xE6,0xF7,0xF4) if ri==4 else (LIGHT if ri%2 else WHITE))
        if ri==0: pr.runs[0].font.color.rgb=WHITE
txt(s, 0.6, 4.7, 6.6, 0.4, "* modelo de produção", size=10, color=GRAY)
pic(s, EVID_DIR/"matriz_confusao.png", 7.45, 2.0, h=2.45)
pic(s, EVID_DIR/"curva_roc.png", 10.3, 2.0, h=2.45)
txt(s, 7.45, 4.6, 5.6, 0.4, "Matriz de confusão e curva ROC (conjunto de teste)", size=11, color=GRAY)
box(s, 0.6, 5.35, 12.15, 1.25, fill=RGBColor(0xE6,0xF7,0xF4))
txt(s, 0.85, 5.5, 11.7, 1.0, [
    "• KNN na mão (NumPy) e scikit-learn ficaram IDÊNTICOS: concordância de 100% nas predições.",
    f"• O MLP teve o melhor F1 ({mm['mlp']['f1']:.3f}) e recall ({mm['mlp']['recall']:.3f}), sendo escolhido como modelo de produção.".replace(".", ",", 2),
], size=14, space=8, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ---- 7. DASHBOARD ----
s = slide(); title(s, "Dashboard de apoio à decisão (Metabase)", "Indicadores, dados tratados, resultados do modelo e fila por risco")
pic(s, EVID_DIR/"dashboard_mockup.png", 2.55, 1.75, h=5.45)

# ---- 8. DECISOES ----
s = slide(); title(s, "Decisões que a solução apoia")
itens=[("Priorizar a fila","Atacar primeiro os chamados com maior probabilidade de violar o SLA."),
       ("Realocar agentes","Direcionar esforço para prioridade Crítica/Alta e categorias mais demoradas."),
       ("Agir na estrutura","Bug/Defeito e Solicitação de Recurso concentram as violações: base de conhecimento e automações."),
       ("Acompanhar no tempo","Monitorar a evolução mensal da taxa de violação e do CSAT.")]
for i,(t,d) in enumerate(itens):
    y=1.9+i*1.2; box(s,0.6,y,12.1,1.05,fill=LIGHT)
    txt(s,0.9,y+0.12,3.4,0.8,t,size=16,bold=True,color=NAVY,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.3,y+0.12,8.1,0.8,d,size=13.5,color=INK,anchor=MSO_ANCHOR.MIDDLE)

# ---- 9. LIMITACOES / CONCLUSAO (dark) ----
s = slide(NAVY)
txt(s, 0.7, 0.6, 12, 0.9, "Limitações e próximos passos", size=30, color=WHITE, bold=True, font="Cambria")
txt(s, 0.7, 1.7, 12, 4.7, [
    "Nuvem implementada e plugável:",
    "• S3 com upload real (boto3) + CloudFormation (IaC), Snowflake (profile dbt), Airflow (DAG) e Metabase.",
    "• Roda 100% local hoje (DuckDB espelha o Snowflake); basta definir as credenciais para ativar a nuvem.",
    "",
    "Limitações: dados sintéticos; integração de nuvem validada localmente (S3 testado com mock).",
    "Próximos passos (escala): Amazon MWAA, SageMaker, embeddings de texto e calibração por custo.",
    "",
    "Conclusão: a solução cobre o ciclo completo, da ingestão à visualização para a decisão,",
    "integrando Modelagem de Dados, Aprendizagem de Máquina e Cloud Computing.",
], size=14, color=RGBColor(0xCA,0xDC,0xFC), space=10)
txt(s, 0.7, 6.6, 12, 0.5, "Obrigado! Dúvidas?", size=18, color=RGBColor(0x02,0xC3,0x9A), bold=True)

out = ROOT/"docs"/"Apresentacao_Projeto_Final.pptx"
prs.save(str(out))
print("apresentacao salva:", out, "| slides:", len(prs.slides._sldIdLst))
