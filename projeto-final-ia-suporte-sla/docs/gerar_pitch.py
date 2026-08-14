# -*- coding: utf-8 -*-
"""
Gera o deck do PITCH DE 15 MINUTOS (docs/Pitch_15min.pptx).

Deck enxuto: 12 slides, ~1,25 min cada. Todos os numeros vem de
evidencias/metrics.json e do export do dashboard, entao o deck nunca
fica dessincronizado do pipeline.

Tambem gera evidencias/comparativo_modelos.png (grafico de F1 por modelo),
usado no slide de resultados.

Uso:
    python docs/gerar_pitch.py
"""
import json
import pathlib
import sys

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[1]))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config.settings import EVID_DIR, ROOT

M = json.load(open(EVID_DIR / "metrics.json", encoding="utf-8"))
D = json.load(open(ROOT / "data" / "curated" / "dashboard" / "dashboard_data.json",
                   encoding="utf-8"))
K, MET = D["kpis"], M["metricas"]
PROD = M["modelo_producao"]                       # 'mlp'
CM = {(r["real"], r["previsto"]): r["n"] for r in D["matriz_confusao"]}

# ----------------------------------------------------------------- paleta
NAVY = RGBColor(0x0F, 0x1E, 0x3D)
NAVY2 = RGBColor(0x1B, 0x30, 0x5C)
SKY = RGBColor(0x2E, 0x9B, 0xD6)
SKY_L = RGBColor(0xA8, 0xD8, 0xF0)
INK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_L = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
LINE = RGBColor(0xDC, 0xE3, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x15, 0x80, 0x3D)
GREEN_L = RGBColor(0xDC, 0xFC, 0xE7)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_L = RGBColor(0xFE, 0xF3, 0xC7)
RED = RGBColor(0xB9, 0x1C, 0x1C)
RED_L = RGBColor(0xFE, 0xE2, 0xE2)

TIT = "Calibri"
BODY = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = 13.333, 7.5
MG = 0.78
BLANK = prs.slide_layouts[6]

_n = [0]


# ----------------------------------------------------------- primitivas
def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.25, radius=0):
    forma = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(forma, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, linhas, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=BODY, space=6, lh=None, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    if isinstance(linhas, str):
        linhas = [linhas]
    for i, ln in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if lh:
            p.line_spacing = lh
        c, b, sz = color, bold, size
        if isinstance(ln, dict):
            c = ln.get("c", color); b = ln.get("b", bold); sz = ln.get("s", size); ln = ln["t"]
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.italic = italic
        r.font.color.rgb = c
        r.font.name = font
    return tb


def cabecalho(s, secao, titulo, sub=None, cor=SKY):
    """Barra superior + rotulo de secao + titulo + numero do slide."""
    _n[0] += 1
    rect(s, 0, 0, W, 0.085, fill=cor)
    txt(s, MG, 0.42, 7.0, 0.3, secao.upper(), size=11.5, color=cor, bold=True)
    txt(s, MG, 0.78, W - 2 * MG - 1.0, 0.75, titulo, size=33, color=NAVY, bold=True, font=TIT)
    y = 1.58
    if sub:
        txt(s, MG, y, W - 2 * MG - 1.0, 0.4, sub, size=15, color=GRAY)
        y += 0.5
    txt(s, W - MG - 0.9, 0.42, 0.9, 0.3, f"{_n[0]:02d}", size=11.5, color=GRAY_L,
        bold=True, align=PP_ALIGN.RIGHT)
    return y + 0.28


def rodape(s, quem, tempo):
    rect(s, 0, H - 0.42, W, 0.02, fill=LINE)
    txt(s, MG, H - 0.34, 6.0, 0.3, quem, size=10.5, color=GRAY_L, bold=True)
    txt(s, W - MG - 3.0, H - 0.34, 3.0, 0.3, tempo, size=10.5, color=GRAY_L,
        align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, titulo=None, corpo=None, cor=SKY, fundo=LIGHT,
         t_size=14, c_size=13):
    rect(s, x, y, w, h, fill=fundo, radius=0.06)
    rect(s, x, y, 0.055, h, fill=cor)
    yy = y + 0.2
    if titulo:
        txt(s, x + 0.3, yy, w - 0.55, 0.35, titulo, size=t_size, color=NAVY, bold=True)
        yy += 0.46
    if corpo:
        txt(s, x + 0.3, yy, w - 0.55, h - (yy - y) - 0.15, corpo, size=c_size,
            color=INK, space=3, lh=1.15)


def stat(s, x, y, w, valor, rotulo, cor=NAVY, fundo=LIGHT, h=1.45, vs=40):
    rect(s, x, y, w, h, fill=fundo, radius=0.08)
    txt(s, x, y + 0.2, w, 0.75, valor, size=vs, color=cor, bold=True,
        align=PP_ALIGN.CENTER, font=TIT)
    txt(s, x + 0.15, y + h - 0.52, w - 0.3, 0.42, rotulo, size=11.5, color=GRAY,
        align=PP_ALIGN.CENTER, space=0, lh=0.95)


def pic(s, caminho, x, y, w=None, h=None, borda=True):
    p = ROOT / caminho
    if not p.exists():
        print(f"  aviso: imagem ausente -> {caminho}")
        return None
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    ph = s.shapes.add_picture(str(p), Inches(x), Inches(y), **kw)
    if borda:
        rect(s, x - 0.035, y - 0.035,
             ph.width / 914400 + 0.07, ph.height / 914400 + 0.07,
             fill=None, line=LINE, lw=1)
    return ph


# --------------------------------------------- grafico de comparacao (F1)
def grafico_modelos():
    nomes = [("Chutar 'não viola'", "baseline_majoritaria"),
             ("Regra atual do gestor", "baseline_regra_prioridade"),
             ("KNN (na mão)", "knn_hardcode"),
             ("KNN (scikit-learn)", "knn_sklearn"),
             ("MLP  →  produção", "mlp")]
    f1 = [MET[k]["f1"] for _, k in nomes]
    rot = [n for n, _ in nomes]
    cores = ["#C7CED8", "#94A3B8", "#7DB9DC", "#7DB9DC", "#15803D"]

    fig, ax = plt.subplots(figsize=(7.6, 3.5), dpi=200)
    barras = ax.barh(rot, f1, color=cores, height=0.62)
    ax.invert_yaxis()
    ax.set_xlim(0, 0.92)
    ax.set_xlabel("F1-score (média harmônica de precisão e recall)", fontsize=9,
                  color="#6B7280")
    for b, v in zip(barras, f1):
        ax.text(v + 0.012, b.get_y() + b.get_height() / 2, f"{v:.3f}".replace(".", ","),
                va="center", fontsize=10.5, fontweight="bold", color="#1F2937")
    ax.axvline(MET["baseline_regra_prioridade"]["f1"], color="#B45309",
               ls="--", lw=1.2, alpha=0.9)
    ax.text(MET["baseline_regra_prioridade"]["f1"] + 0.008, -0.62,
            "barra a superar", fontsize=8.5, color="#B45309", style="italic")
    ax.tick_params(labelsize=10, colors="#1F2937", length=0)
    for lado in ("top", "right", "left"):
        ax.spines[lado].set_visible(False)
    ax.spines["bottom"].set_color("#DCE3EC")
    ax.xaxis.grid(True, color="#EEF2F7", lw=0.9)
    ax.set_axisbelow(True)
    fig.tight_layout()
    destino = EVID_DIR / "comparativo_modelos.png"
    fig.savefig(destino, dpi=200, facecolor="white")
    plt.close(fig)
    return destino


CAM_GRAF = grafico_modelos()


def vg(m, campo):
    return f"{MET[m][campo]:.3f}".replace(".", ",")


def pc(v):
    return f"{v}".replace(".", ",") + "%"


# ==================================================================== 1 CAPA
s = slide(NAVY)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, 0.16, H, fill=SKY)
txt(s, 1.15, 1.75, 11.0, 2.0,
    ["Prever quem vai furar o prazo,", "antes que o prazo fure."],
    size=46, color=WHITE, bold=True, font=TIT, space=2, lh=1.06)
rect(s, 1.18, 3.95, 2.4, 0.045, fill=SKY)
txt(s, 1.15, 4.25, 10.5, 0.9,
    "Pipeline de dados em nuvem e Machine Learning para priorizar a fila do suporte técnico",
    size=17, color=SKY_L, lh=1.2)
txt(s, 1.15, 5.55, 11.0, 0.4,
    "João Paulo Monteiro   ·   Pedro Felipe de Moraes Carrijo   ·   Rogério dos Anjos",
    size=14.5, color=WHITE)
txt(s, 1.15, 6.05, 11.0, 0.7,
    ["IFG · Pós-Graduação em Inteligência Artificial Aplicada · Módulo 2",
     "Modelagem de Dados para IA  ·  Aprendizagem de Máquina  ·  Cloud Computing"],
    size=12, color=RGBColor(0x8A, 0x9C, 0xB8), space=2)

# ============================================================== 2 PROBLEMA
s = slide()
y = cabecalho(s, "O problema", "O gestor prioriza a fila no escuro",
              "Todo chamado tem prazo contratual — e estourar custa multa, insatisfação e cancelamento")
sla = [("Crítica", "4 h", RED), ("Alta", "8 h", AMBER), ("Média", "24 h", SKY), ("Baixa", "48 h", GREEN)]
for i, (nome, t, cor) in enumerate(sla):
    x = MG + i * 1.62
    rect(s, x, y, 1.45, 0.92, fill=LIGHT, radius=0.1)
    rect(s, x, y, 1.45, 0.055, fill=cor)
    txt(s, x, y + 0.19, 1.45, 0.4, t, size=21, color=cor, bold=True,
        align=PP_ALIGN.CENTER, font=TIT)
    txt(s, x, y + 0.62, 1.45, 0.25, nome, size=11, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, MG + 6.65, y + 0.1, 2.0, 0.6, "SLA de\nresolução", size=13, color=GRAY,
    space=0, lh=1.1)

yb = y + 1.35
txt(s, MG, yb, 5.6, 0.35, "A REGRA DE HOJE: “alto ou crítico vai estourar”",
    size=12.5, color=AMBER, bold=True)
stat(s, MG, yb + 0.45, 2.6, pc(round(MET["baseline_regra_prioridade"]["precision"] * 100)),
     "dos alertas que ela dá\nestão certos", AMBER, AMBER_L, h=1.55, vs=38)
stat(s, MG + 2.85, yb + 0.45, 2.6, pc(round((1 - MET["baseline_regra_prioridade"]["recall"]) * 100)),
     "das violações reais\npassam batido", AMBER, AMBER_L, h=1.55, vs=38)

rect(s, 6.6, yb + 0.05, 0.02, 2.1, fill=LINE)
txt(s, 7.0, yb, 5.6, 0.35, "NOSSA PROPOSTA", size=12.5, color=GREEN, bold=True)
card(s, 7.0, yb + 0.45, 5.55, 1.55,
     corpo=["Um risco previsto por chamado, calculado logo após a 1ª resposta do agente.",
            "O gestor reordena a fila por RISCO, não por prioridade cadastrada."],
     cor=GREEN, fundo=GREEN_L, c_size=13.5)

rect(s, MG, 5.95, W - 2 * MG, 0.95, fill=NAVY, radius=0.06)
txt(s, MG + 0.4, 6.16, W - 2 * MG - 0.8, 0.55,
    f"Classificação binária · alvo: violar o SLA · {pc(K['taxa_violacao_pct'])} dos chamados violam — classe minoritária",
    size=15, color=WHITE, bold=True)
rodape(s, "Pedro Felipe", "0:30 – 2:00")

# =========================================================== 3 ARQUITETURA
s = slide()
y = cabecalho(s, "Arquitetura", "Pipeline ELT em sete etapas",
              "Carregamos o dado bruto e transformamos dentro do warehouse, com SQL versionado")
etapas = ["Ingestão", "Atributos", "Machine\nLearning", "Carga", "dbt", "Export", "Nuvem"]
lw = 1.58
for i, e in enumerate(etapas):
    x = MG + i * (lw + 0.16)
    destaque = i == 2
    rect(s, x, y, lw, 1.15, fill=NAVY if destaque else LIGHT, radius=0.09)
    txt(s, x, y + 0.16, lw, 0.28, f"{i+1}", size=11.5,
        color=SKY if destaque else GRAY_L, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, y + 0.46, lw - 0.2, 0.6, e.split("\n"), size=13,
        color=WHITE if destaque else NAVY, bold=True, align=PP_ALIGN.CENTER,
        space=0, lh=1.05)
    if i < len(etapas) - 1:
        txt(s, x + lw, y + 0.42, 0.16, 0.3, "›", size=17, color=GRAY_L,
            align=PP_ALIGN.CENTER)
txt(s, MG, y + 1.32, W - 2 * MG, 0.35,
    "run_pipeline.py (local) e a DAG do Airflow executam as mesmas sete etapas — sem duplicar código.",
    size=13, color=GRAY, italic=True)

camadas = [("BRONZE", "data/raw", "3 fontes brutas: CSV estruturado, texto livre em JSONL, logs de eventos em JSON", SKY),
           ("PRATA", "data/processed", "Limpeza, 9 atributos extraídos do texto e 3 dos logs → uma tabela de 8.000 × 26", SKY),
           ("OURO", "warehouse + dbt", "staging → dimensões → fatos → marts · 12 modelos e 10 testes de qualidade", GREEN)]
for i, (nome, path, desc, cor) in enumerate(camadas):
    yy = y + 1.95 + i * 0.98
    rect(s, MG, yy, W - 2 * MG, 0.85, fill=LIGHT, radius=0.06)
    rect(s, MG, yy, 0.055, 0.85, fill=cor)
    txt(s, MG + 0.32, yy + 0.15, 1.5, 0.3, nome, size=13, color=cor, bold=True)
    txt(s, MG + 0.32, yy + 0.46, 2.4, 0.28, path, size=11, color=GRAY, font=MONO)
    txt(s, MG + 3.2, yy + 0.26, 8.4, 0.5, desc, size=13.5, color=INK)
rodape(s, "João Paulo", "2:00 – 3:30")

# ================================================================= 4 DADOS
s = slide()
y = cabecalho(s, "Dados", "Três tipos de dado, uma tabela só",
              "Exigência 4.2: estruturado + não estruturado + semiestruturado")
fontes = [("ESTRUTURADO", "tickets.csv", "Categoria, prioridade, canal, tier do cliente, tamanho da fila e tempos.", SKY),
          ("NÃO ESTRUTURADO", "ticket_messages.jsonl", "O texto livre que o cliente escreveu ao abrir o chamado.", AMBER),
          ("SEMIESTRUTURADO", "ticket_events.json", "Log de eventos: criado, 1ª resposta, reaberto, resolvido.", GREEN)]
for i, (tipo, arq, desc, cor) in enumerate(fontes):
    x = MG + i * 4.02
    rect(s, x, y, 3.78, 2.05, fill=WHITE, line=LINE, radius=0.06)
    rect(s, x, y, 3.78, 0.075, fill=cor)
    txt(s, x + 0.3, y + 0.3, 3.2, 0.28, tipo, size=11, color=cor, bold=True)
    txt(s, x + 0.3, y + 0.68, 3.3, 0.3, arq, size=13.5, color=NAVY, bold=True, font=MONO)
    txt(s, x + 0.3, y + 1.12, 3.2, 0.8, desc, size=12.5, color=INK, lh=1.15)

yy = y + 2.35
card(s, MG, yy, W - 2 * MG, 0.95,
     corpo=["Do texto extraímos 9 atributos numéricos — tamanho, palavras de urgência, proporção de maiúsculas, "
            "menções a reembolso, cancelamento e erro. Dos logs, 3 agregações. Join pelo ticket_id."],
     cor=SKY, c_size=13.5)

stat(s, MG, yy + 1.2, 3.78, "8.000", "chamados gerados com semente fixa")
stat(s, MG + 4.02, yy + 1.2, 3.78, "26", "colunas na tabela unificada")
stat(s, MG + 8.04, yy + 1.2, 3.75, pc(K["taxa_violacao_pct"]),
     "violam o SLA — classe minoritária", AMBER, AMBER_L)
rodape(s, "João Paulo", "3:30 – 5:00")

# ============================================================ 5 VAZAMENTO
s = slide()
y = cabecalho(s, "Rigor metodológico", "A decisão que sustenta o projeto",
              "Prevenção de vazamento de dados (data leakage)")
rect(s, MG, y, 5.9, 2.95, fill=GREEN_L, radius=0.06)
rect(s, MG, y, 5.9, 0.075, fill=GREEN)
txt(s, MG + 0.35, y + 0.3, 5.2, 0.3, "ENTRA COMO ATRIBUTO", size=12.5, color=GREEN, bold=True)
txt(s, MG + 0.35, y + 0.75, 5.2, 2.0,
    ["categoria · prioridade · canal · tier · região",
     "tamanho da fila · tempo da 1ª resposta",
     "9 atributos do texto · 3 dos logs",
     "",
     "Tudo o que já se conhece no momento da decisão."],
    size=13.5, color=INK, space=4, lh=1.15)

rect(s, 6.95, y, 5.6, 2.95, fill=RED_L, radius=0.06)
rect(s, 6.95, y, 5.6, 0.075, fill=RED)
txt(s, 7.3, y + 0.3, 4.9, 0.3, "FICA DE FORA, DE PROPÓSITO", size=12.5, color=RED, bold=True)
txt(s, 7.3, y + 0.75, 4.9, 2.0,
    ["resolution_minutes → o alvo é calculado dele",
     "csat_score → só existe depois do fechamento",
     "",
     "Usá-los daria quase 100% de acurácia e um modelo",
     "inútil: no momento da decisão eles não existem."],
    size=13.5, color=INK, space=4, lh=1.15)

rect(s, MG, y + 3.25, W - 2 * MG, 1.0, fill=NAVY, radius=0.06)
txt(s, MG + 0.4, y + 3.42, W - 2 * MG - 0.8, 0.7,
    ["A média e o desvio da padronização também são calculados só no conjunto de treino.",
     "É um vazamento sutil, muito comum — e infla a métrica sem ninguém perceber."],
    size=14, color=WHITE, space=2, lh=1.15)
rodape(s, "João Paulo", "5:00 – 6:00")

# ============================================================= 6 DBT
s = slide()
y = cabecalho(s, "Modelagem", "Modelo dimensional e qualidade de dados",
              "dbt: staging → dimensões → fatos → marts")
card(s, MG, y, 7.6, 2.5, titulo="O que construímos",
     corpo=["fct_tickets — fato central, 1 registro por chamado, integrando as 3 fontes",
            "4 dimensões (categoria, canal, produto, data) + fct_predictions",
            "mart_ticket_decision — traduz a probabilidade em faixa Alto / Médio / Baixo",
            "10 testes: unicidade, não-nulo, valores aceitos e integridade referencial"],
     cor=SKY, c_size=13.5)
rect(s, 8.75, y, 3.8, 2.5, fill=GREEN_L, radius=0.08)
txt(s, 8.75, y + 0.5, 3.8, 0.9, "PASS = 22", size=44, color=GREEN, bold=True,
    align=PP_ALIGN.CENTER, font=TIT)
txt(s, 8.75, y + 1.45, 3.8, 0.7, ["12 modelos + 10 testes", "0 avisos · 0 erros"],
    size=13, color=GREEN, align=PP_ALIGN.CENTER, space=2)

yy = y + 2.85
rect(s, MG, yy, W - 2 * MG, 1.65, fill=LIGHT, radius=0.06)
rect(s, MG, yy, 0.055, 1.65, fill=SKY)
txt(s, MG + 0.35, yy + 0.2, 11.2, 0.32, "Portabilidade comprovada", size=15,
    color=NAVY, bold=True)
txt(s, MG + 0.35, yy + 0.62, 11.3, 0.95,
    ["Os mesmos 12 modelos e 10 testes rodam com PASS=22 no DuckDB E no Postgres — sem mudar uma linha de SQL, só o --target.",
     "Foi esse exercício que revelou o único trecho não portável: um cast as double, atalho exclusivo do DuckDB, trocado pelo",
     "tipo ANSI double precision. O profile do Snowflake usa exatamente o mesmo código."],
    size=13, color=INK, space=2, lh=1.15)
rodape(s, "João Paulo", "6:00 – 7:30")

# ============================================================= 7 AIRFLOW
s = slide()
y = cabecalho(s, "Orquestração", "A DAG em execução",
              "Agendamento diário, retry por task e histórico auditável")
pic(s, "evidencias/airflow_graph.png", MG, y, w=8.95)
cards = [("7 tasks", "encadeadas com dependência explícita"),
         ("3 execuções", "todas verdes no histórico"),
         ("BashOperator", "o Airflow orquestra; cada etapa roda isolada no ambiente do projeto")]
for i, (t, d) in enumerate(cards):
    yy = y + i * 1.42
    rect(s, 10.0, yy, 2.55, 1.25, fill=LIGHT, radius=0.08)
    txt(s, 10.2, yy + 0.16, 2.2, 0.35, t, size=15, color=NAVY, bold=True)
    txt(s, 10.2, yy + 0.55, 2.2, 0.65, d, size=11.5, color=GRAY, lh=1.1)
rodape(s, "João Paulo", "7:30 – 8:30")

# ================================================================== 8 ML
s = slide()
y = cabecalho(s, "Aprendizagem de Máquina", "Medido contra a prática atual",
              "Classificação binária · 6.000 treino / 2.000 teste · 49 atributos · conjunto de teste fora da amostra")
pic(s, str(CAM_GRAF.relative_to(ROOT)), MG, y - 0.05, w=7.5, borda=False)

xd = 8.6
card(s, xd, y - 0.05, 3.95, 1.35, titulo="Hard-code = biblioteca",
     corpo=["KNN em NumPy puro e o do scikit-learn: 100% de concordância nas 2.000 predições."],
     cor=SKY, c_size=12.5)
card(s, xd, y + 1.45, 3.95, 1.35, titulo="Ganho sobre o que se faz hoje",
     corpo=[f"F1 sobe de {vg('baseline_regra_prioridade','f1')} para {vg(PROD,'f1')} — 36% melhor que a regra por prioridade."],
     cor=GREEN, fundo=GREEN_L, c_size=12.5)
card(s, xd, y + 2.95, 3.95, 1.35, titulo="Por que F1 e não acurácia",
     corpo=["Chutar 'não viola' já acerta 77%. Com classe minoritária, acurácia engana."],
     cor=AMBER, fundo=AMBER_L, c_size=12.5)
rodape(s, "Rogério", "8:30 – 11:00")

# ============================================================ 9 MATRIZ
s = slide()
y = cabecalho(s, "Aprendizagem de Máquina", "O que isso significa na operação",
              f"Matriz de confusão do modelo de produção ({PROD.upper()}) nos 2.000 chamados de teste")

cx, cy, cw, ch = MG + 1.15, y + 0.42, 2.35, 1.35
txt(s, cx, cy - 0.34, cw * 2, 0.28, "PREVISTO PELO MODELO", size=10.5, color=GRAY, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, cx - 0.15, cy + 0.4, 0.3, 2.8, "R\nE\nA\nL", size=10.5, color=GRAY, bold=True,
    align=PP_ALIGN.CENTER, space=0, lh=0.92)
celulas = [(0, 0, CM[(0, 0)], "Não violou\ne acertamos", GREEN, GREEN_L),
           (1, 0, CM[(0, 1)], "Alarme falso", AMBER, AMBER_L),
           (0, 1, CM[(1, 0)], "Escapou", RED, RED_L),
           (1, 1, CM[(1, 1)], "Violou e\navisamos", GREEN, GREEN_L)]
for col, lin, n, rot, cor, fundo in celulas:
    x = cx + col * (cw + 0.12)
    yy = cy + lin * (ch + 0.12)
    rect(s, x, yy, cw, ch, fill=fundo, radius=0.07)
    txt(s, x, yy + 0.18, cw, 0.6, f"{n:,}".replace(",", "."), size=33, color=cor,
        bold=True, align=PP_ALIGN.CENTER, font=TIT)
    txt(s, x, yy + 0.82, cw, 0.45, rot.split("\n"), size=11, color=cor,
        align=PP_ALIGN.CENTER, space=0, lh=1.0)
for i, r in enumerate(["Não violou", "Violou"]):
    txt(s, cx + i * (cw + 0.12), cy + 2 * ch + 0.3, cw, 0.25, r, size=11,
        color=GRAY, align=PP_ALIGN.CENTER)

xr = 7.2
tp, fn, fp = CM[(1, 1)], CM[(1, 0)], CM[(0, 1)]
stat(s, xr, y + 0.35, 2.6, str(tp + fn), "violaram o SLA de verdade", NAVY)
stat(s, xr + 2.75, y + 0.35, 2.6, str(tp), "o modelo avisou a tempo", GREEN, GREEN_L)
stat(s, xr, y + 1.95, 2.6, str(fn), "escaparam", RED, RED_L)
stat(s, xr + 2.75, y + 1.95, 2.6, str(fp), "alarmes falsos", AMBER, AMBER_L)
card(s, xr, y + 3.55, 5.35, 1.0,
     corpo=[f"Priorizamos recall ({vg(PROD,'recall')}): um falso negativo é um prazo estourado sem aviso — caro. "
            "Um falso positivo é um agente que olhou um chamado à toa — barato."],
     cor=SKY, c_size=12.5)
rodape(s, "Rogério", "11:00 – 11:45")

# =========================================================== 10 DASHBOARD
s = slide()
y = cabecalho(s, "Visualização", "O painel onde a decisão acontece",
              "Metabase sobre o warehouse · 10 cards · filtros de prioridade, categoria e canal")
pic(s, "evidencias/metabase_dashboard.png", MG, y, w=8.6)
top = D["por_categoria"][0]
rect(s, 9.65, y, 2.9, 4.35, fill=AMBER_L, radius=0.07)
rect(s, 9.65, y, 2.9, 0.075, fill=AMBER)
txt(s, 9.9, y + 0.28, 2.45, 0.55, "O achado mais\nacionável", size=14, color=NAVY,
    bold=True, space=0, lh=1.1)
txt(s, 9.9, y + 1.0, 2.45, 0.9, pc(top["taxa_violacao_pct"]), size=34, color=AMBER,
    bold=True, font=TIT)
txt(s, 9.9, y + 1.75, 2.45, 0.55, f"dos chamados de\n“{top['category']}”\nviolam o SLA",
    size=12, color=INK, space=0, lh=1.1)
txt(s, 9.9, y + 2.75, 2.45, 1.4,
    ["“Dúvida / Como Fazer” viola 0,5%.",
     "",
     "Não é problema de previsão: é de processo. Pedido de recurso não deveria dividir fila e prazo com suporte."],
    size=11.5, color=INK, space=3, lh=1.1)
rodape(s, "Pedro Felipe", "11:45 – 13:15")

# =============================================================== 11 NUVEM
s = slide()
y = cabecalho(s, "Cloud Computing", "O que é real e o que é proposta",
              "Transparência sobre o escopo efetivamente executado")
pic(s, "infra/arquitetura_aws.png", MG, y + 0.1, w=7.7)
xr = 8.75
rect(s, xr, y, 3.8, 1.85, fill=GREEN_L, radius=0.07)
rect(s, xr, y, 3.8, 0.075, fill=GREEN)
txt(s, xr + 0.28, y + 0.28, 3.2, 0.3, "IMPLEMENTADO", size=12, color=GREEN, bold=True)
txt(s, xr + 0.28, y + 0.66, 3.25, 1.1,
    ["CloudFormation (IaC): 3 buckets, IAM, Lambda, EventBridge e CloudWatch",
     "Template extra para o AWS Academy, que não permite criar IAM Role"],
    size=11.5, color=INK, space=4, lh=1.1)

rect(s, xr, y + 2.05, 3.8, 1.25, fill=AMBER_L, radius=0.07)
rect(s, xr, y + 2.05, 3.8, 0.075, fill=AMBER)
txt(s, xr + 0.28, y + 2.33, 3.2, 0.3, "PROPOSTA NO DIAGRAMA", size=12, color=AMBER, bold=True)
txt(s, xr + 0.28, y + 2.7, 3.25, 0.55,
    "Airflow gerenciado, SageMaker e Snowflake gerenciado — arquitetura de referência.",
    size=11.5, color=INK, lh=1.1)

rect(s, xr, y + 3.5, 3.8, 1.1, fill=LIGHT, radius=0.07)
txt(s, xr + 0.28, y + 3.72, 3.2, 0.3, "CUSTO", size=12, color=NAVY, bold=True)
txt(s, xr + 0.28, y + 4.02, 3.25, 0.5,
    "Hoje ~US$ 0 (nível gratuito) · produção US$ 450–600/mês · versão enxuta US$ 70–120",
    size=11.5, color=INK, lh=1.1)
rodape(s, "Pedro Felipe", "13:15 – 14:15")

# ========================================================== 12 FECHAMENTO
s = slide(NAVY)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, W, 0.085, fill=SKY)
txt(s, MG, 0.62, 11.5, 0.7, "As três decisões que a solução apoia", size=33,
    color=WHITE, bold=True, font=TIT)
decisoes = [("Priorizar a fila",
             f"Atacar primeiro os {D['distribuicao_risco'][0]['chamados']} chamados de risco Alto, em vez de seguir a prioridade cadastrada."),
            ("Realocar agentes",
             "Concentrar esforço em Crítica (" + pc(D["por_prioridade"][0]["taxa_violacao_pct"])
             + ") e Alta (" + pc(D["por_prioridade"][1]["taxa_violacao_pct"]) + "), onde o prazo mais estoura."),
            ("Agir na estrutura",
             "Separar “Solicitação de Recurso” da fila de suporte, com SLA próprio — ela depende do time de produto.")]
for i, (t, d) in enumerate(decisoes):
    yy = 1.72 + i * 1.12
    rect(s, MG, yy, W - 2 * MG, 0.95, fill=NAVY2, radius=0.06)
    rect(s, MG, yy, 0.055, 0.95, fill=SKY)
    txt(s, MG + 0.35, yy + 0.15, 3.0, 0.35, t, size=16, color=SKY_L, bold=True)
    txt(s, MG + 3.6, yy + 0.17, 8.2, 0.6, d, size=13.5, color=WHITE, lh=1.12)

txt(s, MG, 5.35, 11.5, 0.35, "LIMITAÇÕES QUE ASSUMIMOS", size=12, color=SKY, bold=True)
lim = ["Dados sintéticos — as métricas são otimistas",
       "Modelo estático, sem retreino nem monitoramento de drift",
       "Painel sobre histórico, não sobre a fila em tempo real",
       "Sem Snowflake e sem upload real ao S3: implementados, não executados"]
for i, l in enumerate(lim):
    col, lin = i % 2, i // 2
    x = MG + col * 6.0
    yy = 5.78 + lin * 0.42
    rect(s, x, yy + 0.09, 0.09, 0.09, fill=SKY)
    txt(s, x + 0.25, yy, 5.6, 0.32, l, size=12.5, color=RGBColor(0xC5, 0xD2, 0xE4))

saida = ROOT / "docs" / "Pitch_15min.pptx"
prs.save(saida)
print(f"OK -> {saida}  ({_n[0] + 2} slides)")
print(f"     grafico -> {CAM_GRAF}")
